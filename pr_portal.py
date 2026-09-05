"""
ЦСР PR-ПОРТАЛ 3.0
Премиальный AI-ассистент для пресс-службы ЦСР
- Генерация постов с GPT-4o
- Мониторинг новостей (News API + Google RSS)
- Интеграция с Яндекс.Диском
- Аналитика и метрики
- История постов

Требования: streamlit>=1.52 (нужен callable в st.download_button),
openai, requests, feedparser, beautifulsoup4, python-docx, python-pptx,
pandas, openpyxl, pypdf.
"""

import base64
import hashlib
import io
import json
import os
import re
import uuid
from email.utils import parsedate_to_datetime
from pathlib import Path
from datetime import datetime, timedelta, timezone
import requests
import feedparser
import urllib.parse

import docx
import pandas as pd
import streamlit as st
from bs4 import BeautifulSoup
from openai import OpenAI
from pptx import Presentation

# pypdf — актуальное продолжение заброшенного PyPDF2.
# Fallback оставлен, чтобы приложение не падало на старом окружении.
try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    from PyPDF2 import PdfReader

# ============================================================
# ФИРМЕННАЯ ТЕМА ЦСР
# ============================================================

CSR_THEME_CSS = """
<style>
/* ============================================================
   ЦСР PR-ПОРТАЛ — PREMIUM 3D EDITORIAL UI
   Фирменная эстетика: glassmorphism + layered 3D + soft glow
   ============================================================ */
:root {
    --csr-blue: #221AC8;
    --csr-indigo: #34306F;
    --csr-violet: #5E5B94;
    --csr-lilac: #8E8AC6;
    --csr-lavender: #C7C5F3;
    --csr-milk: #F7F6FF;
    --csr-text: #2B285F;
    --csr-text-soft: #514D86;
    --csr-text-muted: #7773A9;
    --csr-green: #55A98E;
    --csr-cyan: #5D9FC8;
    --csr-orange: #D99862;
    --csr-red: #C9777A;

    --glass-1: rgba(255,255,255,.48);
    --glass-2: rgba(255,255,255,.30);
    --glass-3: rgba(234,232,255,.34);
    --glass-edge: rgba(255,255,255,.74);
    --glass-edge-soft: rgba(255,255,255,.44);
    --blue-edge: rgba(34,26,200,.18);

    --shadow-xl: 0 28px 70px rgba(47,42,111,.19);
    --shadow-lg: 0 18px 44px rgba(47,42,111,.15);
    --shadow-md: 0 10px 28px rgba(47,42,111,.11);
    --inner-light: inset 0 1px 0 rgba(255,255,255,.72);
    --glow-blue: 0 0 30px rgba(63,54,220,.12);

    --radius-xl: 30px;
    --radius-lg: 23px;
    --radius-md: 17px;
    --radius-sm: 13px;
}

* { box-sizing: border-box; }
html, body { background: #C9C8F1; }
html, body, [data-testid="stAppViewContainer"] { color: var(--csr-text); }

.stApp {
    background:
        radial-gradient(circle at 7% 4%, rgba(255,255,255,.82) 0 8%, rgba(255,255,255,0) 27%),
        radial-gradient(circle at 91% 10%, rgba(118,112,224,.22) 0 8%, rgba(118,112,224,0) 30%),
        radial-gradient(circle at 74% 93%, rgba(255,255,255,.34) 0 7%, rgba(255,255,255,0) 25%),
        linear-gradient(145deg, #E0DFFF 0%, #CFCEF8 34%, #BAB9EC 100%);
    background-attachment: fixed;
    min-height: 100vh;
}

/* subtle editorial grid */
.stApp::before {
    content: "";
    position: fixed;
    inset: 0;
    pointer-events: none;
    opacity: .13;
    background-image:
        linear-gradient(rgba(255,255,255,.45) 1px, transparent 1px),
        linear-gradient(90deg, rgba(255,255,255,.45) 1px, transparent 1px);
    background-size: 44px 44px;
    mask-image: linear-gradient(to bottom, rgba(0,0,0,.45), transparent 72%);
    z-index: 0;
}

[data-testid="stAppViewContainer"] > .main { position: relative; z-index: 1; }
.block-container {
    max-width: 1500px;
    padding-top: 1.55rem !important;
    padding-bottom: 2.5rem !important;
    padding-left: 2.25rem !important;
    padding-right: 2.25rem !important;
}

/* ---------- Typography ---------- */
h1, h2, h3, h4, h5, h6 {
    color: var(--csr-blue) !important;
    font-weight: 800 !important;
    letter-spacing: -.025em !important;
}
h1 { font-size: 2.15rem !important; line-height: 1.05 !important; }
h2 { font-size: 1.34rem !important; }
h3 { font-size: 1.08rem !important; }
[data-testid="stMarkdownContainer"] p {
    color: var(--csr-text-soft) !important;
    line-height: 1.62 !important;
}
[data-testid="stCaptionContainer"], .stCaption, small {
    color: var(--csr-text-muted) !important;
}

/* ---------- Sidebar ---------- */
[data-testid="stSidebar"] {
    background:
        radial-gradient(circle at 12% 5%, rgba(255,255,255,.70), transparent 24%),
        linear-gradient(180deg, rgba(246,245,255,.56), rgba(223,221,255,.34)) !important;
    backdrop-filter: blur(28px) saturate(130%);
    -webkit-backdrop-filter: blur(28px) saturate(130%);
    border-right: 1px solid rgba(255,255,255,.65);
    box-shadow: 14px 0 34px rgba(62,56,139,.08);
}
[data-testid="stSidebar"] .block-container {
    padding: 1.25rem 1.05rem 1.5rem !important;
}
[data-testid="stSidebar"] * { color: var(--csr-text) !important; }

.sidebar-brand {
    position: relative;
    overflow: hidden;
    padding: 17px 17px 15px;
    margin: 0 0 14px;
    border-radius: 25px;
    background:
        radial-gradient(circle at 12% 12%, rgba(255,255,255,.78), transparent 38%),
        linear-gradient(145deg, rgba(255,255,255,.52), rgba(232,230,255,.29));
    border: 1px solid var(--glass-edge);
    box-shadow: var(--shadow-md), var(--inner-light), var(--glow-blue);
    backdrop-filter: blur(22px);
}
.sidebar-brand::after {
    content: "";
    position: absolute;
    width: 110px; height: 110px;
    border-radius: 36px;
    right: -43px; top: -48px;
    transform: rotate(26deg);
    background: linear-gradient(145deg, rgba(34,26,200,.16), rgba(255,255,255,.04));
    border: 1px solid rgba(255,255,255,.30);
}
.sidebar-logo-wrap {
    position: relative;
    z-index: 2;
    width: 188px;
    max-width: 92%;
    min-height: 57px;
    display: flex;
    align-items: center;
    margin-bottom: 8px;
}
.sidebar-logo-wrap img {
    display: block;
    width: 100%;
    height: auto;
    object-fit: contain;
    filter: drop-shadow(0 5px 10px rgba(39,34,94,.08));
}
.sidebar-brand-fallback {
    font-weight: 900;
    color: var(--csr-blue) !important;
    font-size: 1.52rem;
    letter-spacing: -.04em;
}
.sidebar-subtitle {
    position: relative;
    z-index: 2;
    max-width: 220px;
    color: var(--csr-text-soft) !important;
    font-size: .83rem;
    line-height: 1.45;
}
.sidebar-section-label {
    margin: 15px 3px 8px;
    color: var(--csr-text-muted) !important;
    font-size: .66rem;
    font-weight: 800;
    text-transform: uppercase;
    letter-spacing: .13em;
}
.storage-pill {
    display: inline-flex;
    align-items: center;
    gap: 7px;
    margin: 1px 0 3px;
    padding: 7px 10px;
    border: 1px solid rgba(255,255,255,.54);
    border-radius: 999px;
    background: rgba(255,255,255,.28);
    box-shadow: var(--inner-light);
    font-size: .72rem;
    color: var(--csr-text-soft) !important;
}
.storage-dot { width: 7px; height: 7px; border-radius: 50%; display: inline-block; }
.storage-ok { background: var(--csr-green); box-shadow: 0 0 10px rgba(85,169,142,.45); }
.storage-warn { background: var(--csr-orange); box-shadow: 0 0 10px rgba(217,152,98,.42); }
.storage-bad { background: var(--csr-red); box-shadow: 0 0 10px rgba(201,119,122,.42); }

/* ---------- Hero / 3D editorial scene ---------- */
.csr-hero {
    position: relative;
    overflow: hidden;
    min-height: 220px;
    margin: 0 0 1.25rem;
    padding: 30px 32px;
    border-radius: 32px;
    border: 1px solid rgba(255,255,255,.76);
    background:
        radial-gradient(circle at 10% 12%, rgba(255,255,255,.76), transparent 34%),
        radial-gradient(circle at 88% 17%, rgba(70,62,199,.18), transparent 31%),
        linear-gradient(145deg, rgba(255,255,255,.45), rgba(223,220,255,.24));
    box-shadow: var(--shadow-xl), var(--inner-light), var(--glow-blue);
    backdrop-filter: blur(26px) saturate(130%);
    -webkit-backdrop-filter: blur(26px) saturate(130%);
}
.csr-hero::before {
    content: "";
    position: absolute;
    width: 390px; height: 390px;
    right: -180px; top: -210px;
    border-radius: 50%;
    border: 1px solid rgba(255,255,255,.38);
    box-shadow: inset 0 0 55px rgba(255,255,255,.14), 0 0 55px rgba(34,26,200,.09);
}
.hero-layout {
    position: relative;
    z-index: 2;
    display: grid;
    grid-template-columns: minmax(0, 1.2fr) minmax(330px, .8fr);
    gap: 28px;
    align-items: center;
}
.hero-kicker {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    margin-bottom: 12px;
    padding: 7px 12px;
    border: 1px solid rgba(255,255,255,.58);
    border-radius: 999px;
    background: rgba(255,255,255,.30);
    box-shadow: var(--inner-light);
    color: var(--csr-blue) !important;
    font-size: .72rem;
    font-weight: 800;
    letter-spacing: .08em;
    text-transform: uppercase;
}
.hero-title {
    margin: 0 0 10px;
    color: var(--csr-blue) !important;
    font-size: clamp(2.25rem, 4vw, 3.7rem);
    line-height: .98;
    font-weight: 900;
    letter-spacing: -.055em;
}
.hero-subtitle {
    max-width: 760px;
    color: var(--csr-text-soft) !important;
    font-size: 1rem;
    line-height: 1.58;
}
.hero-signals {
    display: flex;
    gap: 8px;
    margin-top: 18px;
    flex-wrap: wrap;
}
.hero-signal {
    padding: 7px 11px;
    border-radius: 999px;
    background: rgba(255,255,255,.27);
    border: 1px solid rgba(255,255,255,.48);
    color: var(--csr-text-soft) !important;
    font-size: .72rem;
    font-weight: 700;
}
.hero-scene {
    position: relative;
    height: 165px;
    perspective: 900px;
    transform-style: preserve-3d;
}
.scene-core {
    position: absolute;
    width: 128px; height: 128px;
    left: 50%; top: 50%;
    transform: translate(-50%,-50%) rotateX(59deg) rotateZ(-34deg);
    border-radius: 29px;
    background: linear-gradient(145deg, rgba(255,255,255,.67), rgba(111,104,209,.24));
    border: 1px solid rgba(255,255,255,.82);
    box-shadow: 22px 27px 38px rgba(48,42,117,.20), inset 0 2px 0 rgba(255,255,255,.84), 0 0 30px rgba(59,50,212,.12);
}
.scene-core::before {
    content: "AI";
    position: absolute;
    inset: 0;
    display: grid;
    place-items: center;
    color: var(--csr-blue);
    font-size: 2rem;
    font-weight: 900;
    transform: translateZ(1px);
}
.scene-panel {
    position: absolute;
    width: 105px; height: 62px;
    border-radius: 17px;
    border: 1px solid rgba(255,255,255,.72);
    background: linear-gradient(150deg, rgba(255,255,255,.51), rgba(206,203,248,.18));
    box-shadow: 12px 17px 28px rgba(50,45,117,.16), inset 0 1px 0 rgba(255,255,255,.78);
    backdrop-filter: blur(12px);
}
.scene-panel::after {
    content: "";
    position: absolute;
    left: 14px; right: 14px; bottom: 14px;
    height: 5px;
    border-radius: 999px;
    background: linear-gradient(90deg, var(--csr-blue) 0 42%, rgba(255,255,255,.38) 42% 100%);
    box-shadow: 0 0 12px rgba(34,26,200,.18);
}
.sp1 { left: 2%; top: 8%; transform: rotate(-8deg) translateZ(65px); }
.sp2 { right: 1%; top: 18%; transform: rotate(8deg) translateZ(25px); }
.sp3 { right: 12%; bottom: 1%; width: 86px; height: 52px; transform: rotate(-5deg) translateZ(80px); }
.scene-orb {
    position: absolute;
    border-radius: 50%;
    background: radial-gradient(circle at 30% 24%, rgba(255,255,255,.95), rgba(126,118,218,.40) 42%, rgba(34,26,200,.24) 74%);
    border: 1px solid rgba(255,255,255,.75);
    box-shadow: 0 12px 24px rgba(49,43,118,.16), inset 0 1px 0 rgba(255,255,255,.75), 0 0 18px rgba(34,26,200,.12);
}
.so1 { width: 35px; height: 35px; left: 13%; bottom: 12%; }
.so2 { width: 22px; height: 22px; right: 20%; top: 0; }
.scene-line {
    position: absolute;
    height: 1px;
    background: linear-gradient(90deg, transparent, rgba(34,26,200,.36), transparent);
    box-shadow: 0 0 9px rgba(34,26,200,.18);
    transform-origin: left center;
}
.sl1 { width: 115px; left: 19%; top: 66%; transform: rotate(-26deg); }
.sl2 { width: 125px; left: 53%; top: 24%; transform: rotate(24deg); }

/* ---------- Section header ---------- */
.section-head { margin: 0 0 .65rem; padding: 0 2px; }
.section-eyebrow {
    color: var(--csr-text-muted) !important;
    text-transform: uppercase;
    letter-spacing: .13em;
    font-size: .66rem;
    font-weight: 800;
    margin-bottom: 4px;
}
.section-title {
    color: var(--csr-blue) !important;
    font-size: 1.22rem;
    font-weight: 850;
    letter-spacing: -.025em;
    line-height: 1.1;
}
.section-note {
    color: var(--csr-text-soft) !important;
    font-size: .82rem;
    margin-top: 4px;
}

/* ---------- Glass containers ---------- */
[data-testid="stVerticalBlockBorderWrapper"] {
    background:
        radial-gradient(circle at 12% 4%, rgba(255,255,255,.42), transparent 36%),
        linear-gradient(160deg, rgba(255,255,255,.36), rgba(224,221,255,.22));
    border: 1px solid rgba(255,255,255,.67) !important;
    border-radius: var(--radius-lg) !important;
    box-shadow: var(--shadow-lg), var(--inner-light), 0 0 26px rgba(60,52,212,.06);
    backdrop-filter: blur(24px) saturate(125%);
    -webkit-backdrop-filter: blur(24px) saturate(125%);
}

/* ---------- Metrics ---------- */
.metric-card {
    position: relative;
    overflow: hidden;
    min-height: 145px;
    padding: 20px 19px 17px;
    border-radius: 22px;
    border: 1px solid rgba(255,255,255,.70);
    background:
        radial-gradient(circle at 90% 8%, rgba(255,255,255,.53), transparent 30%),
        linear-gradient(155deg, rgba(255,255,255,.42), rgba(223,220,255,.22));
    box-shadow: var(--shadow-md), var(--inner-light), var(--glow-blue);
    backdrop-filter: blur(20px);
    transition: transform .22s ease, box-shadow .22s ease;
}
.metric-card:hover {
    transform: translateY(-3px);
    box-shadow: 0 19px 38px rgba(47,42,111,.16), var(--inner-light), var(--glow-blue);
}
.metric-orb {
    position: absolute;
    right: 16px; top: 15px;
    width: 38px; height: 38px;
    border-radius: 13px;
    background: linear-gradient(145deg, rgba(255,255,255,.58), rgba(108,102,203,.20));
    border: 1px solid rgba(255,255,255,.72);
    box-shadow: 7px 9px 17px rgba(49,43,118,.13), inset 0 1px 0 rgba(255,255,255,.80);
}
.metric-label {
    position: relative;
    z-index: 2;
    max-width: calc(100% - 48px);
    color: var(--csr-text-muted) !important;
    font-size: .69rem;
    font-weight: 800;
    text-transform: uppercase;
    letter-spacing: .105em;
}
.metric-value {
    position: relative;
    z-index: 2;
    margin: 16px 0 7px;
    color: var(--csr-blue) !important;
    font-size: 2.55rem;
    line-height: .94;
    font-weight: 900;
    letter-spacing: -.035em;
}
.metric-note {
    position: relative;
    z-index: 2;
    color: var(--csr-text-muted) !important;
    font-size: .75rem;
    line-height: 1.35;
}

/* ---------- Inputs ---------- */
input, textarea, select,
[data-baseweb="input"], [data-baseweb="textarea"], [data-baseweb="select"] > div {
    background: linear-gradient(180deg, rgba(255,255,255,.52), rgba(244,243,255,.34)) !important;
    color: var(--csr-text) !important;
    border: 1px solid rgba(255,255,255,.70) !important;
    border-radius: 14px !important;
    box-shadow: var(--inner-light), 0 7px 18px rgba(47,42,111,.07) !important;
    backdrop-filter: blur(15px);
}
input:focus, textarea:focus {
    border-color: rgba(34,26,200,.34) !important;
    box-shadow: 0 0 0 3px rgba(34,26,200,.09), var(--inner-light) !important;
}
input::placeholder, textarea::placeholder { color: #8A87B6 !important; }

/* Select slider */
[data-testid="stSlider"] [role="slider"] {
    background: var(--csr-blue) !important;
    box-shadow: 0 0 0 5px rgba(34,26,200,.10), 0 4px 12px rgba(34,26,200,.18) !important;
}

/* ---------- Buttons ---------- */
.stButton > button, .stDownloadButton > button {
    min-height: 42px !important;
    border-radius: 14px !important;
    border: 1px solid rgba(255,255,255,.70) !important;
    background: linear-gradient(180deg, rgba(255,255,255,.54), rgba(235,232,255,.35)) !important;
    color: var(--csr-blue) !important;
    font-weight: 750 !important;
    box-shadow: var(--shadow-md), var(--inner-light) !important;
    transition: transform .18s ease, box-shadow .18s ease, background .18s ease !important;
}
.stButton > button:hover, .stDownloadButton > button:hover {
    transform: translateY(-1px);
    background: linear-gradient(180deg, rgba(255,255,255,.70), rgba(239,237,255,.43)) !important;
    box-shadow: 0 14px 27px rgba(47,42,111,.14), var(--inner-light) !important;
}
button[kind="primary"] {
    background: linear-gradient(135deg, #221AC8 0%, #4D47D8 52%, #6660B4 100%) !important;
    color: white !important;
    border: 1px solid rgba(255,255,255,.25) !important;
    box-shadow: 0 13px 28px rgba(34,26,200,.28), inset 0 1px 0 rgba(255,255,255,.28), 0 0 22px rgba(34,26,200,.10) !important;
}
button[kind="primary"]:hover {
    box-shadow: 0 17px 34px rgba(34,26,200,.34), inset 0 1px 0 rgba(255,255,255,.32) !important;
}

/* ---------- Tabs ---------- */
[role="tablist"] {
    gap: 7px;
    padding: 4px 4px 8px;
    border-bottom: 1px solid rgba(255,255,255,.38);
}
[role="tab"] {
    min-height: 38px;
    padding: 7px 12px !important;
    border-radius: 12px !important;
    color: var(--csr-text-muted) !important;
    font-weight: 700 !important;
}
[role="tab"]:hover { background: rgba(255,255,255,.28) !important; }
[role="tab"][aria-selected="true"] {
    color: var(--csr-blue) !important;
    background: linear-gradient(180deg, rgba(255,255,255,.52), rgba(237,235,255,.34)) !important;
    border: 1px solid rgba(255,255,255,.55) !important;
    box-shadow: var(--inner-light), 0 7px 17px rgba(47,42,111,.07);
}

/* ---------- Expanders / alerts / uploader / code ---------- */
[data-testid="stAlert"],
[data-testid="stExpander"],
[data-testid="stFileUploaderDropzone"],
[data-testid="stCodeBlock"], pre {
    background:
        radial-gradient(circle at 10% 2%, rgba(255,255,255,.47), transparent 36%),
        linear-gradient(155deg, rgba(255,255,255,.37), rgba(229,227,255,.24)) !important;
    border: 1px solid rgba(255,255,255,.66) !important;
    border-radius: 18px !important;
    box-shadow: var(--shadow-md), var(--inner-light) !important;
    backdrop-filter: blur(19px);
}
[data-testid="stFileUploaderDropzone"] { border-style: dashed !important; }
[data-testid="stExpander"] summary { color: var(--csr-blue) !important; font-weight: 750 !important; }
[data-testid="stCodeBlock"] code, pre code, code {
    color: var(--csr-text) !important;
    background: transparent !important;
    font-size: .91rem !important;
}

/* ---------- Charts ---------- */
[data-testid="stVegaLiteChart"], [data-testid="stArrowVegaLiteChart"] {
    background: rgba(255,255,255,.20);
    border: 1px solid rgba(255,255,255,.45);
    border-radius: 18px;
    padding: 8px;
    box-shadow: var(--inner-light);
}

/* ---------- Misc ---------- */
a { color: var(--csr-blue) !important; font-weight: 700; text-decoration: none !important; }
a:hover { text-decoration: underline !important; }
hr {
    border: none !important;
    height: 1px !important;
    background: linear-gradient(90deg, transparent, rgba(255,255,255,.86), transparent) !important;
    margin: 1.05rem 0 !important;
}
[data-testid="stMetricValue"] { color: var(--csr-blue) !important; font-weight: 900 !important; }
footer, #MainMenu { visibility: hidden; }

/* ---------- Responsive ---------- */
@media (max-width: 980px) {
    .hero-layout { grid-template-columns: 1fr; }
    .hero-scene { display: none; }
    .csr-hero { min-height: 0; padding: 24px; }
    .block-container { padding-left: 1.1rem !important; padding-right: 1.1rem !important; }
}
</style>
"""

# ============================================================
# КОНФИГ
# ============================================================
st.set_page_config(
    page_title="ЦСР PR-портал 3.0", 
    page_icon="⚡️", 
    layout="wide", 
    initial_sidebar_state="expanded"
)

st.markdown(CSR_THEME_CSS, unsafe_allow_html=True)

# ============================================================
# БРЕНДИНГ / ЛОГОТИП
# ============================================================
APP_DIR = Path(__file__).resolve().parent


def find_logo_path():
    """Надёжно находит логотип независимо от cwd Streamlit Cloud."""
    candidates = [
        APP_DIR / "csr_logo.png",
        APP_DIR / "csr_logo(1).png",
        APP_DIR / "assets" / "csr_logo.png",
        Path.cwd() / "csr_logo.png",
        Path.cwd() / "csr_logo(1).png",
    ]
    for candidate in candidates:
        if candidate.exists() and candidate.is_file():
            return candidate

    # Последний fallback: любой csr_logo*.png рядом с приложением.
    dynamic = sorted(APP_DIR.glob("csr_logo*.png"))
    return dynamic[0] if dynamic else None


@st.cache_data(show_spinner=False)
def image_as_data_uri(path_str):
    try:
        path = Path(path_str)
        payload = base64.b64encode(path.read_bytes()).decode("ascii")
        return f"data:image/png;base64,{payload}"
    except Exception:
        return ""


def logo_uri():
    path = find_logo_path()
    return image_as_data_uri(str(path)) if path else ""


def render_sidebar_brand():
    uri = logo_uri()
    if uri:
        logo_html = f'<img src="{uri}" alt="ЦСР">'
    else:
        logo_html = '<div class="sidebar-brand-fallback">ЦСР</div>'

    st.markdown(
        f"""
        <div class="sidebar-brand">
            <div class="sidebar-logo-wrap">{logo_html}</div>
            <div class="sidebar-subtitle">PR-портал · единая цифровая среда пресс-службы</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_storage_status():
    if storage is None:
        status_class = "storage-warn"
        label = "Локальный режим хранения"
    elif storage.writable is False:
        status_class = "storage-bad"
        label = "Яндекс.Диск: ошибка записи"
    else:
        status_class = "storage-ok"
        label = "Яндекс.Диск подключён"
    st.markdown(
        f'<div class="storage-pill"><span class="storage-dot {status_class}"></span>{label}</div>',
        unsafe_allow_html=True,
    )


def render_main_hero():
    st.markdown(
        """
        <div class="csr-hero">
          <div class="hero-layout">
            <div>
              <div class="hero-kicker">ЦСР · AI workspace</div>
              <div class="hero-title">PR-портал 3.0</div>
              <div class="hero-subtitle">
                Интеллектуальная среда для подготовки публикаций, проверки фактуры,
                мониторинга инфополя, работы с базой знаний и аналитики контента.
              </div>
              <div class="hero-signals">
                <span class="hero-signal">AI-редактор</span>
                <span class="hero-signal">База знаний</span>
                <span class="hero-signal">Мониторинг</span>
                <span class="hero-signal">Аналитика</span>
              </div>
            </div>
            <div class="hero-scene" aria-hidden="true">
              <div class="scene-line sl1"></div>
              <div class="scene-line sl2"></div>
              <div class="scene-panel sp1"></div>
              <div class="scene-panel sp2"></div>
              <div class="scene-panel sp3"></div>
              <div class="scene-core"></div>
              <div class="scene-orb so1"></div>
              <div class="scene-orb so2"></div>
            </div>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_section_header(eyebrow, title, note=""):
    note_html = f'<div class="section-note">{note}</div>' if note else ""
    st.markdown(
        f"""
        <div class="section-head">
            <div class="section-eyebrow">{eyebrow}</div>
            <div class="section-title">{title}</div>
            {note_html}
        </div>
        """,
        unsafe_allow_html=True,
    )

MODEL = "gpt-4o"
TEMP_GENERATE = 0.4
TEMP_REFINE = 0.3
TEMP_FACTCHECK = 0.1

STYLE_LIBRARY_FILE = "style_library.json"
MAX_STYLE_EXAMPLES_IN_PROMPT = 5
MAX_CHARS_PER_EXAMPLE = 1200
HISTORY_KEEP_LAST = 6

# Путь к папке базы знаний на Яндекс.Диске.
# Если папка называется иначе — поменяйте здесь одну строку.
# Точное имя можно посмотреть в разделе «📁» (кнопка «Показать корень диска»).
YANDEX_BASE_PATH = "/ЦСР - база знаний"

CATEGORY_FIELDS = {
    "1. Анонс мероприятия": [
        {"label": "Дата мероприятия", "key": "date", "area": False},
        {"label": "Название мероприятия", "key": "event_name", "area": False},
        {"label": "Участники / спикеры", "key": "participants", "area": True},
        {"label": "Темы для обсуждения", "key": "topics", "area": True},
        {"label": "Место проведения", "key": "location", "area": False},
    ],
    "2. Исследование": [
        {"label": "Название исследования", "key": "research_name", "area": False},
        {"label": "Источник / автор", "key": "source", "area": False},
        {"label": "Ключевые цифры и выводы", "key": "key_findings", "area": True},
        {"label": "Методология (кратко, необязательно)", "key": "methodology", "area": True},
    ],
    "3. Пост после мероприятия (Итоги)": [
        {"label": "Название мероприятия", "key": "event_name", "area": False},
        {"label": "Дата", "key": "date", "area": False},
        {"label": "Спикеры", "key": "participants", "area": True},
        {"label": "Основные итоги / тезисы", "key": "outcomes", "area": True},
    ],
    "4. Публикация в СМИ (Колонка / комментарий)": [
        {"label": "Название СМИ", "key": "media_name", "area": False},
        {"label": "Автор / спикер", "key": "author", "area": False},
        {"label": "Тема комментария", "key": "topic", "area": True},
        {"label": "Ссылка на публикацию", "key": "link", "area": False},
    ],
    "5. Тематический пост": [
        {"label": "Тема", "key": "topic", "area": False},
        {"label": "Ключевые тезисы", "key": "key_points", "area": True},
        {"label": "Источник (если есть)", "key": "source", "area": False},
    ],
    "6. Индексы": [
        {"label": "Название индекса", "key": "index_name", "area": False},
        {"label": "Значение / показатель", "key": "value", "area": False},
        {"label": "Динамика (рост/падение)", "key": "dynamics", "area": False},
        {"label": "Источник", "key": "source", "area": False},
    ],
    "7. Публикация каталогов": [
        {"label": "Название каталога", "key": "catalog_name", "area": False},
        {"label": "Ссылка / файл", "key": "link", "area": False},
        {"label": "Краткое описание содержания", "key": "description", "area": True},
    ],
    "8. Еженедельная статистика (Коротко о ценах)": [
        {"label": "Период", "key": "period", "area": False},
        {"label": "Категории и цифры", "key": "figures", "area": True},
        {"label": "Источник данных", "key": "source", "area": False},
    ],
    "9. Рубрика: Цифры и факты": [
        {"label": "Список фактов (цифра — пояснение, по одному на строку)", "key": "facts_list", "area": True},
        {"label": "Источники (по каждому факту)", "key": "sources", "area": True},
    ],
}

# ============================================================
# СЕКРЕТЫ И КЛЮЧИ
# ============================================================
SECRETS_AVAILABLE = True


def get_secret(name, default=None):
    """Безопасно читает secrets. Отсутствие файла секретов фиксируется
    отдельным флагом: молча продолжать в этом случае нельзя — так когда-то
    незаметно отключалась парольная защита."""
    global SECRETS_AVAILABLE
    try:
        return st.secrets.get(name, default)
    except Exception:
        SECRETS_AVAILABLE = False
        return default


APP_PASSWORD = get_secret("APP_PASSWORD")
API_KEY = get_secret("OPENAI_API_KEY")
YANDEX_TOKEN = get_secret("YANDEX_DISK_TOKEN")
NEWS_API_KEY = get_secret("NEWS_API_KEY")

if APP_PASSWORD:
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if not st.session_state.authenticated:
        st.title("ЦСР PR-портал")
        pwd = st.text_input("Введите пароль доступа", type="password")
        if st.button("Войти", type="primary"):
            if pwd == APP_PASSWORD:
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("❌ Неверный пароль")
        st.stop()

if not API_KEY:
    st.error("⚠️ Не найден OPENAI_API_KEY в secrets.")
    if not SECRETS_AVAILABLE:
        st.info(
            "Файл секретов вообще не читается. Создайте `.streamlit/secrets.toml` "
            "локально или заполните Secrets в настройках Streamlit Cloud."
        )
    st.stop()

client = OpenAI(api_key=API_KEY)

# ============================================================
# ЯНДЕКС.ДИСК API
# ============================================================
class YandexDiskAPI:
    def __init__(self, token):
        self.token = token
        self.headers = {"Authorization": f"OAuth {token}"}
        self.base_url = "https://cloud-api.yandex.net/v1/disk"
    
    def list_files(self, path="/"):
        url = f"{self.base_url}/resources"
        params = {"path": path, "limit": 100}
        try:
            response = requests.get(url, headers=self.headers, params=params, timeout=10)
            if response.status_code == 200:
                items = response.json().get("_embedded", {}).get("items", [])
                return sorted(items, key=lambda x: x.get("name") or "")
            return []
        except Exception:
            return []

    def get_download_url(self, path):
        url = f"{self.base_url}/resources/download"
        params = {"path": path}
        try:
            response = requests.get(url, headers=self.headers, params=params, timeout=10)
            if response.status_code == 200:
                return response.json().get("href")
        except Exception:
            pass
        return None

    def download_file(self, path):
        download_url = self.get_download_url(path)
        if download_url:
            try:
                response = requests.get(download_url, timeout=30)
                if response.status_code == 200:
                    return response.content
            except Exception:
                pass
        return None

disk_api = None
if YANDEX_TOKEN:
    disk_api = YandexDiskAPI(YANDEX_TOKEN)


# ============================================================
# ПОСТОЯННОЕ ХРАНИЛИЩЕ (Яндекс.Диск)
# ------------------------------------------------------------
# Файловая система Streamlit Cloud временная: при перезапуске
# приложения всё, что лежит рядом с кодом, стирается. Поэтому
# история постов, избранные новости и библиотека стилей хранятся
# в JSON-файлах на Яндекс.Диске, а локально держится только кэш.
# ============================================================

# Папка для служебных данных приложения. Создаётся автоматически.
APP_DATA_PATH = "/PR-портал ЦСР — данные"


class StorageUnavailable(Exception):
    """Хранилище временно недоступно. Важно отличать это от «файла ещё нет»:
    в первом случае писать нельзя, иначе пустой список затрёт все данные."""


class YandexStorage:
    def __init__(self, token, folder=APP_DATA_PATH):
        self.token = token
        self.folder = folder
        self.headers = {"Authorization": f"OAuth {token}"}
        self.base_url = "https://cloud-api.yandex.net/v1/disk"
        self.writable = None      # None = ещё не проверяли
        self.last_error = ""

    def probe(self):
        """Разовая проверка токена и прав на запись при старте приложения.
        Без неё индикатор в сайдбаре горел зелёным даже с битым токеном."""
        if self.writable is None:
            self.writable = self._ensure_folder()
        return self.writable

    def _ensure_folder(self):
        """Создаёт папку приложения. 409 = уже существует, это норма."""
        try:
            r = requests.put(
                f"{self.base_url}/resources",
                headers=self.headers,
                params={"path": self.folder},
                timeout=10,
            )
            if r.status_code in (201, 409):
                return True
            self.last_error = f"{r.status_code}: {_api_message(r)}"
            return False
        except Exception as e:
            self.last_error = str(e)
            return False

    def load(self, name, default, strict=False):
        """Читает JSON с Диска.

        strict=False — «мягкое» чтение для отображения: при любой проблеме
        отдаёт default.
        strict=True  — чтение перед записью: 404 (файла ещё нет) отдаёт
        default, а сетевая ошибка поднимает StorageUnavailable, чтобы
        вызывающий код отменил запись и не потерял данные.
        """
        if not self.token:
            if strict:
                raise StorageUnavailable("Токен Яндекс.Диска не задан")
            return default
        try:
            r = requests.get(
                f"{self.base_url}/resources/download",
                headers=self.headers,
                params={"path": f"{self.folder}/{name}"},
                timeout=10,
            )
            if r.status_code == 404:
                return default          # файла ещё нет — первый запуск
            if r.status_code != 200:
                self.last_error = f"{r.status_code}: {_api_message(r)}"
                if strict:
                    raise StorageUnavailable(self.last_error)
                return default
            href = r.json().get("href")
            data = requests.get(href, timeout=15)
            if data.status_code != 200:
                self.last_error = f"download {data.status_code}"
                if strict:
                    raise StorageUnavailable(self.last_error)
                return default
            return json.loads(data.content.decode("utf-8"))
        except StorageUnavailable:
            raise
        except Exception as e:
            self.last_error = str(e)
            if strict:
                raise StorageUnavailable(str(e))
            return default

    def save(self, name, data):
        """Записывает JSON на Диск. Возвращает True при успехе."""
        if not self.token:
            self.writable = False
            return False
        try:
            self._ensure_folder()
            r = requests.get(
                f"{self.base_url}/resources/upload",
                headers=self.headers,
                params={"path": f"{self.folder}/{name}", "overwrite": "true"},
                timeout=10,
            )
            if r.status_code != 200:
                self.writable = False
                self.last_error = f"{r.status_code}: {_api_message(r)}"
                return False
            href = r.json().get("href")
            payload = json.dumps(data, ensure_ascii=False, indent=2).encode("utf-8")
            up = requests.put(href, data=payload, timeout=30)
            ok = up.status_code in (201, 202)
            self.writable = ok
            if not ok:
                self.last_error = f"upload {up.status_code}"
            return ok
        except Exception as e:
            self.writable = False
            self.last_error = str(e)
            return False


def _api_message(response):
    """Достаёт человекочитаемое сообщение об ошибке из ответа API Диска."""
    try:
        return response.json().get("message", "")
    except Exception:
        return response.text[:200]


storage = YandexStorage(YANDEX_TOKEN) if YANDEX_TOKEN else None


def _local_load(name, default):
    if os.path.exists(name):
        try:
            with open(name, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return default
    return default


def storage_load(name, default):
    """Мягкое чтение для отображения. Никогда не бросает исключений."""
    if storage:
        return storage.load(name, default)
    return _local_load(name, default)


def storage_load_for_write(name, default):
    """Чтение перед записью. Бросает StorageUnavailable, если Диск не отвечает,
    чтобы вызывающий код НЕ перезаписал удалённый файл неполными данными."""
    if storage:
        return storage.load(name, default, strict=True)
    return _local_load(name, default)


def storage_save(name, data):
    """Пишет данные на Диск и одновременно в локальный файл (быстрый кэш)."""
    try:
        with open(name, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass
    if storage:
        return storage.save(name, data)
    return True  # локальный режим: запись в файл считается успешной


def guarded_write(file_name, mutate, default):
    """Общий безопасный сценарий «прочитать → изменить → записать».

    Возвращает True при успехе. При недоступном хранилище показывает ошибку
    и НИЧЕГО не пишет — это защита от потери истории при сетевом сбое.
    """
    try:
        current = storage_load_for_write(file_name, default)
    except StorageUnavailable as e:
        st.error(
            f"⚠️ Хранилище недоступно, изменения не сохранены ({e}). "
            "Попробуйте ещё раз через минуту."
        )
        return False
    if not isinstance(current, type(default)):
        current = default
    updated = mutate(current)
    if updated is None:
        return True  # менять нечего
    if not storage_save(file_name, updated):
        st.error(
            "⚠️ Не удалось записать данные на Яндекс.Диск"
            + (f": {storage.last_error}" if storage and storage.last_error else "")
        )
        return False
    return True

# ============================================================
# МОНИТОРИНГ НОВОСТЕЙ
# ============================================================

def parse_any_date(value):
    """Приводит дату к datetime с таймзоной.

    News API отдаёт ISO 8601 (2026-01-15T09:30:00Z), Google News RSS —
    RFC 822 (Wed, 15 Jan 2026 09:30:00 GMT). Раньше их сортировали как
    строки, из-за чего порядок при двух включённых источниках был случайным.
    """
    if not value:
        return None
    text = str(value).strip()
    try:
        dt = datetime.fromisoformat(text.replace("Z", "+00:00"))
        return dt if dt.tzinfo else dt.replace(tzinfo=timezone.utc)
    except (ValueError, TypeError):
        pass
    try:
        dt = parsedate_to_datetime(text)
        return dt if dt.tzinfo else dt.replace(tzinfo=timezone.utc)
    except (TypeError, ValueError, IndexError):
        return None


def format_news_date(news_or_value):
    """Единый человекочитаемый формат даты для любого источника."""
    if isinstance(news_or_value, dict):
        dt = news_or_value.get("published_dt") or parse_any_date(news_or_value.get("published"))
    else:
        dt = parse_any_date(news_or_value)
    return dt.strftime("%d.%m.%Y") if dt else "дата неизвестна"


def widget_key(prefix, *parts):
    """Стабильный уникальный ключ виджета.

    Обрезка заголовка до 20 символов давала коллизии на однотипных
    новостных заголовках и падение с DuplicateElementKey.
    """
    raw = "|".join(str(p) for p in parts)
    digest = hashlib.md5(raw.encode("utf-8")).hexdigest()[:12]
    return f"{prefix}_{digest}"


class NewsMonitor:
    def __init__(self, news_api_key=None):
        self.news_api_key = news_api_key
        self.news_api_url = "https://newsapi.org/v2/everything"
        self.init_db()
    
    FAVORITES_FILE = "news_favorites.json"

    def init_db(self):
        pass  # хранилище файловое, инициализация не нужна

    def save_favorite(self, title, source, url, published, description):
        def mutate(items):
            if any(i.get("title") == title for i in items):
                return None       # уже в избранном — записывать нечего
            items.insert(0, {
                "title": title,
                "source": source,
                "url": url,
                "published": published,
                "description": description,
                "saved_at": datetime.now().isoformat(timespec="seconds"),
            })
            return items

        return guarded_write(self.FAVORITES_FILE, mutate, [])

    def get_favorites(self):
        items = storage_load(self.FAVORITES_FILE, [])
        if not isinstance(items, list):
            return []
        return [
            {
                "title": i.get("title") or "",
                "source": i.get("source") or "",
                "url": i.get("url") or "",
                "published": i.get("published") or "",
                "description": i.get("description") or "",
            }
            for i in items
        ]

    def delete_favorite(self, title):
        def mutate(items):
            return [i for i in items if i.get("title") != title]

        return guarded_write(self.FAVORITES_FILE, mutate, [])

    def fetch_news_api(self, query="ЦСР", language="ru", days=7):
        if not self.news_api_key:
            return []
        try:
            params = {
                "q": query,
                "language": language,
                "sortBy": "publishedAt",
                "apiKey": self.news_api_key,
                "pageSize": 25,
                "from": (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d"),
            }
            response = requests.get(self.news_api_url, params=params, timeout=10)
            if response.status_code != 200:
                st.warning(
                    f"News API вернул {response.status_code}: {_api_message(response)}"
                )
                return []
            data = response.json()
            news_list = []
            for article in data.get("articles") or []:
                # ВАЖНО: .get(key, default) отдаёт None, если ключ есть,
                # но значение null — а News API регулярно так делает.
                # Поэтому везде `or ""`, иначе дальше падает published[:10].
                source_obj = article.get("source") or {}
                news_list.append({
                    "title": (article.get("title") or "").strip() or "Без заголовка",
                    "source": (source_obj.get("name") or "News API").strip(),
                    "url": article.get("url") or "",
                    "published": article.get("publishedAt") or "",
                    "published_dt": parse_any_date(article.get("publishedAt")),
                    "description": article.get("description") or "",
                    "type": "news_api",
                })
            return news_list
        except Exception as e:
            st.warning(f"News API недоступен: {e}")
            return []

    def fetch_google_news(self, query="ЦСР OR \"Центр стратегических разработок\"", days=7):
        try:
            # ВАЖНО: запрос обязательно кодируем, иначе пробелы ломают URL
            q = f"{query} when:{days}d"
            url = (
                "https://news.google.com/rss/search?q="
                + urllib.parse.quote(q)
                + "&hl=ru&gl=RU&ceid=RU:ru"
            )
            feed = feedparser.parse(url)
            news_list = []
            for entry in feed.entries[:40]:
                try:
                    title = entry.get("title") or ""
                    link = entry.get("link") or ""
                    published = entry.get("published") or ""
                    summary = entry.get("summary") or ""
                    if summary:
                        summary = BeautifulSoup(summary, "html.parser").get_text()[:250]
                    source_obj = entry.get("source") or {}
                    source = source_obj.get("title") if source_obj else None
                    if not source and " - " in title:
                        title, source = title.rsplit(" - ", 1)
                    news_list.append({
                        "title": title.strip() or "Без заголовка",
                        "source": (source or "Google News").strip(),
                        "url": link,
                        "published": published,
                        "published_dt": parse_any_date(published),
                        "description": summary,
                        "type": "google_news",
                    })
                except Exception:
                    continue
            return news_list
        except Exception as e:
            st.warning(f"Google News недоступен: {e}")
            return []

    def fetch_combined(self, query="ЦСР", use_news_api=True, use_google_news=True, days=7):
        all_news = []
        if use_news_api and self.news_api_key:
            all_news.extend(self.fetch_news_api(query, days=days))
        if use_google_news:
            all_news.extend(self.fetch_google_news(query, days=days))

        seen_titles = set()
        unique_news = []
        for news in all_news:
            title = news.get("title") or ""
            key = title.strip().lower()
            if key and key not in seen_titles:
                seen_titles.add(key)
                unique_news.append(news)

        # Сортируем по datetime, а не по строке: News API отдаёт ISO,
        # Google News — RFC 822, и лексикографическое сравнение их путает.
        far_past = datetime.min.replace(tzinfo=timezone.utc)
        return sorted(
            unique_news,
            key=lambda x: x.get("published_dt") or far_past,
            reverse=True,
        )

news_monitor = NewsMonitor(news_api_key=NEWS_API_KEY)


# Кэш на 10 минут: без него новости перезапрашиваются при КАЖДОМ клике
# в интерфейсе и дневная квота News API (100 запросов) сгорает за пару минут.
#
# ВНИМАНИЕ: параметр НЕ должен начинаться с подчёркивания. st.cache_data
# намеренно исключает такие аргументы из ключа кэша, поэтому прежний
# `_cache_buster` не сбрасывал кэш и кнопка «Обновить» ничего не делала.
@st.cache_data(ttl=600, show_spinner=False)
def cached_fetch_news(query, use_api, use_google, days, cache_buster=0):
    return news_monitor.fetch_combined(
        query, use_news_api=use_api, use_google_news=use_google, days=days
    )


@st.cache_data(ttl=300, show_spinner=False)
def cached_list_disk(path, cache_buster=0):
    if not disk_api:
        return []
    return disk_api.list_files(path)


def disk_buster():
    """Счётчик для принудительного обновления листинга Яндекс.Диска."""
    return st.session_state.get("disk_cache_buster", 0)


def bump_disk_cache():
    st.session_state.disk_cache_buster = disk_buster() + 1


def _streamlit_at_least(major, minor):
    try:
        parts = st.__version__.split(".")
        return (int(parts[0]), int(parts[1])) >= (major, minor)
    except (AttributeError, IndexError, ValueError):
        return False


# Передача callable в data= у st.download_button появилась только в 1.52.
# На более старых версиях это падало с «Invalid binary data format: <class
# 'function'>», поэтому здесь есть двухшаговый запасной вариант.
SUPPORTS_DEFERRED_DOWNLOAD = _streamlit_at_least(1, 52)


FILE_ICONS = {
    "pdf": "📄",
    "txt": "📝", "docx": "📝", "doc": "📝", "html": "📝", "htm": "📝",
    "xlsx": "📊", "xls": "📊", "csv": "📊",
    "pptx": "🎬", "ppt": "🎬",
    "ai": "🎨", "psd": "🎨", "png": "🎨", "jpg": "🎨", "jpeg": "🎨",
}


def file_icon(file_name):
    ext = (file_name or "").lower().rsplit(".", 1)[-1]
    return FILE_ICONS.get(ext, "📑")


def format_size(size_bytes):
    """КБ для мелких файлов: раньше всё меньше 50 КБ показывалось как «0.0 MB»."""
    size_bytes = size_bytes or 0
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.0f} КБ"
    return f"{size_bytes / (1024 * 1024):.1f} МБ"


def render_disk_file_row(file_info, key_prefix):
    file_name = file_info.get("name") or ""
    col1, col2, col3 = st.columns([2.5, 0.7, 0.6])
    with col1:
        st.write(f"{file_icon(file_name)} {file_name}")
    with col2:
        st.caption(format_size(file_info.get("size")))
    with col3:
        disk_download_button(file_info, key_prefix)


def disk_download_button(file_info, key_prefix):
    """Кнопка скачивания файла с Яндекс.Диска.

    Файл подгружается только по нажатию — иначе при открытии базы знаний
    приложение скачивало бы содержимое всех папок сразу.
    """
    path = file_info.get("path")
    name = file_info.get("name") or "file"
    base_key = widget_key(key_prefix, path)

    if SUPPORTS_DEFERRED_DOWNLOAD:
        # Внутри callable команды Streamlit игнорируются, поэтому сообщить
        # об ошибке нельзя: при неудаче отдаём пустой файл.
        st.download_button(
            "📥",
            data=lambda p=path: (disk_api.download_file(p) or b""),
            file_name=name,
            key=base_key,
            use_container_width=True,
        )
        return

    buffer_key = f"{base_key}__buffer"
    if st.session_state.get(buffer_key):
        st.download_button(
            "💾",
            data=st.session_state[buffer_key],
            file_name=name,
            key=f"{base_key}__save",
            use_container_width=True,
        )
    elif st.button("📥", key=f"{base_key}__fetch", use_container_width=True):
        with st.spinner("Скачиваю..."):
            content = disk_api.download_file(path)
        if content:
            st.session_state[buffer_key] = content
            st.rerun()
        else:
            st.error(f"Не удалось скачать «{name}»")

# ============================================================
# ИСТОРИЯ ПОСТОВ (хранится на Яндекс.Диске)
# ============================================================
POSTS_FILE = "posts_history.json"


# Имя параметра снова без подчёркивания — иначе сброс кэша не работает
# и только что сохранённый пост не появляется в истории и метриках.
@st.cache_data(ttl=60, show_spinner=False)
def _load_posts(buster=0):
    data = storage_load(POSTS_FILE, [])
    return data if isinstance(data, list) else []


def _bump_posts_cache():
    st.session_state.posts_buster = st.session_state.get("posts_buster", 0) + 1


def save_post(category, tone, length, text, raw_source, has_citations=0):
    """Дописывает пост в общий журнал. Перед записью перечитывает файл,
    чтобы не затереть посты, добавленные коллегой с другого компьютера.
    Если Диск недоступен — запись отменяется, а не перетирает историю."""
    post_id = str(uuid.uuid4())

    def mutate(posts):
        posts.append({
            "id": post_id,
            "category": category,
            "tone": tone,
            "length": length,
            "text": text,
            "raw_source": (raw_source or "")[:4000],
            "has_citations": int(has_citations),
            "created_at": datetime.now().isoformat(timespec="seconds"),
        })
        return posts

    if not guarded_write(POSTS_FILE, mutate, []):
        return None
    _bump_posts_cache()
    return post_id


def get_posts_history(limit=100):
    posts = _load_posts(st.session_state.get("posts_buster", 0))
    posts = sorted(posts, key=lambda p: p.get("created_at", ""), reverse=True)
    return [
        (p.get("id"), p.get("category"), p.get("tone"),
         p.get("text", ""), p.get("created_at", ""), p.get("has_citations", 0))
        for p in posts[:limit]
    ]


def get_metrics():
    posts = _load_posts(st.session_state.get("posts_buster", 0))
    now = datetime.now()

    def parse(p):
        try:
            return datetime.fromisoformat(p.get("created_at") or "")
        except (TypeError, ValueError):
            return None

    week_ago = now - timedelta(days=7)
    month_ago = now - timedelta(days=30)
    dates = [d for d in (parse(p) for p in posts) if d]

    def group(field):
        counts = {}
        for p in posts:
            key = p.get(field) or "—"
            counts[key] = counts.get(key, 0) + 1
        return sorted(counts.items(), key=lambda x: x[1], reverse=True)

    with_citations = sum(1 for p in posts if p.get("has_citations"))

    return {
        "total_posts": len(posts),
        "posts_week": sum(1 for d in dates if d > week_ago),
        "posts_month": sum(1 for d in dates if d > month_ago),
        "by_category": group("category"),
        "by_tone": group("tone"),
        "by_length": group("length"),
        "with_citations": with_citations,
        "without_citations": len(posts) - with_citations,
    }

# ============================================================
# ФУНКЦИИ
# ============================================================

# Единый список форматов: раньше загрузчик принимал html/htm, а парсер их
# не знал — файлы молча превращались в пустую строку. Теперь один источник
# правды и для file_uploader, и для браузера Яндекс.Диска.
SUPPORTED_EXTENSIONS = ("pdf", "docx", "pptx", "xlsx", "txt", "html", "htm")
STYLE_EXTENSIONS = ("txt", "pdf", "docx", "html", "htm")


def _parse_bytes(content, name):
    """Извлекает текст из байтов по расширению имени. Бросает исключения —
    вызывающий код решает, показывать ли их пользователю."""
    lowered = name.lower()
    text = ""

    if lowered.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    elif lowered.endswith(".docx"):
        doc = docx.Document(io.BytesIO(content))
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Таблицы раньше терялись целиком, хотя в них обычно и лежат цифры.
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    text += " | ".join(cells) + "\n"
    elif lowered.endswith(".pptx"):
        prs = Presentation(io.BytesIO(content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                    text += shape.text_frame.text + "\n"
    elif lowered.endswith(".xlsx"):
        sheets = pd.read_excel(io.BytesIO(content), sheet_name=None)
        for sheet_name, df in sheets.items():
            text += f"--- Лист: {sheet_name} ---\n"
            text += df.to_string(index=False) + "\n\n"
    elif lowered.endswith((".html", ".htm")):
        soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n", strip=True)
    elif lowered.endswith(".txt"):
        text = content.decode("utf-8", errors="ignore")
    else:
        raise ValueError(f"формат «{Path(name).suffix or '?'}» не поддерживается")

    return text


def extract_text_from_file(uploaded_file):
    """Читает файл из st.file_uploader. При ошибке показывает предупреждение,
    а не молча возвращает пустую строку."""
    try:
        uploaded_file.seek(0)
        content = uploaded_file.read()
        return _parse_bytes(content, uploaded_file.name)
    except Exception as e:
        st.warning(f"⚠️ Не удалось прочитать «{uploaded_file.name}»: {e}")
        return ""


def extract_text_from_files(uploaded_files):
    if not uploaded_files:
        return ""
    chunks = []
    for f in uploaded_files:
        t = extract_text_from_file(f)
        if t.strip():
            chunks.append(f"### {f.name}\n{t}")
    return "\n\n".join(chunks)


def extract_yandex_file(disk_client, file_path):
    content = disk_client.download_file(file_path)
    if not content:
        st.warning(f"⚠️ Не удалось скачать {file_path} с Яндекс.Диска")
        return ""
    try:
        return _parse_bytes(content, file_path)
    except Exception as e:
        st.warning(f"⚠️ Не удалось прочитать {file_path}: {e}")
        return ""

# Кэшируем: без этого docx пересобирался при каждой перерисовке вкладки
# «Экспорт», то есть на любой клик в интерфейсе.
@st.cache_data(show_spinner=False)
def create_docx(text):
    doc = docx.Document()
    for line in text.split("\n"):
        if line.strip() == "":
            doc.add_paragraph("")
            continue
        p = doc.add_paragraph()
        parts = re.split(r"(\*\*.+?\*\*)", line)
        for part in parts:
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                run = p.add_run(part[2:-2])
                run.bold = True
            elif part:
                p.add_run(part)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def call_openai(messages, temperature):
    try:
        response = client.chat.completions.create(
            model=MODEL, 
            messages=messages, 
            temperature=temperature
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"❌ Ошибка OpenAI: {e}")
        return None

def trim_history(history):
    """Оставляет system + первое задание + хвост диалога.

    Хвост берём чётной длины, чтобы не разрывать пары user/assistant:
    прежняя версия могла выбросить первый ответ модели и оставить два
    user-сообщения подряд.
    """
    if len(history) <= HISTORY_KEEP_LAST + 2:
        return history
    keep = HISTORY_KEEP_LAST - (HISTORY_KEEP_LAST % 2)
    tail = history[2:][-keep:] if keep else []
    return history[:2] + tail

def load_style_library():
    lib = storage_load(STYLE_LIBRARY_FILE, [])
    return lib if isinstance(lib, list) else []


def add_style_examples(texts):
    def mutate(library):
        for t in texts:
            t = (t or "").strip()
            if not t:
                continue
            library.append({
                "id": str(uuid.uuid4()),
                "preview": (t[:70] + "…") if len(t) > 70 else t,
                "text": t,
                "added": datetime.now().strftime("%Y-%m-%d %H:%M"),
            })
        return library

    return guarded_write(STYLE_LIBRARY_FILE, mutate, [])


def delete_style_example(example_id):
    # .get вместо e["id"]: записи старого формата без ключа роняли приложение
    def mutate(library):
        return [e for e in library if e.get("id") != example_id]

    return guarded_write(STYLE_LIBRARY_FILE, mutate, [])


def build_style_block(selected_ids):
    if not selected_ids:
        return ""
    library = load_style_library()
    chosen = [e for e in library if e.get("id") in selected_ids][:MAX_STYLE_EXAMPLES_IN_PROMPT]
    if not chosen:
        return ""
    parts = [(e.get("text") or "")[:MAX_CHARS_PER_EXAMPLE] for e in chosen]
    block = "\n\n---\n\n".join(p for p in parts if p)
    if not block:
        return ""
    return f"Примеры постов для ориентира:\n\n{block}"

def get_enhanced_system_rules(category, tone, text_length, use_citations):
    citation_rule = "Пересказывай информацию своими словами без цитат" if not use_citations else "Используй цитаты из источников (максимум 2-3 на пост)"
    
    return f"""Ты — PR-редактор ЦСР. Основные правила:

🎯 СТИЛЬ:
• Аналитически и структурированно
• Заголовок первой строкой **жирным**
• Логические переходы между абзацами
• Четкая логика и причинно-следственные связи

📝 КОНТЕНТ:
• БЕЗ буллитов — только прозаический текст
• {citation_rule}
• Конкретные цифры с контекстом
• Завершающий вывод
• Язык живой, но серьёзный

🚫 ИЗБЕГАЙ:
• "и т.д.", "и т.п."
• Клише типа "как известно"
• Излишних эмодзи
• Коротких рубленых предложений

✅ ДЕЛАЙ:
• Связки: "Это означает...", "Поэтому..."
• Анализ, а не перечисление
• Подвал: "⭐️ подписаться на канал"

Рубрика: {category}
Тон: {tone}
Объем: {text_length}
"""

FIELD_KEY_PREFIX = "field_"


def field_widget_key(category, field_key):
    """Ключ виджета включает рубрику.

    Раньше ключ был просто field_<key>, а `date`, `source`, `event_name`
    встречаются в нескольких рубриках — значения перетекали при переключении.
    """
    return widget_key(FIELD_KEY_PREFIX, category, field_key)


def build_dynamic_fields(category):
    values = {}
    fields = CATEGORY_FIELDS.get(category, [])
    for field in fields:
        key = field_widget_key(category, field["key"])
        if field["area"]:
            values[field["key"]] = st.text_area(field["label"], key=key)
        else:
            values[field["key"]] = st.text_input(field["label"], key=key)
    return values


def reset_input_form():
    """Полная очистка формы для кнопки «Новый пост».

    Прежняя версия сбрасывала только значения из `defaults`, а введённый
    текст жил в session_state под ключами виджетов и переживал сброс.
    """
    for key, value in defaults.items():
        st.session_state[key] = value
    st.session_state.disk_content = ""
    st.session_state.news_for_post = None

    stale = [
        k for k in list(st.session_state.keys())
        if k.startswith(FIELD_KEY_PREFIX)
        or k.startswith("disk_")
        or k in ("speaker_input", "materials_input", "disk_folder_select")
    ]
    for k in stale:
        del st.session_state[k]

    # Загрузчики файлов очищаются только сменой ключа
    st.session_state.uploader_nonce = st.session_state.get("uploader_nonce", 0) + 1

def format_fields_as_text(category, values):
    fields = CATEGORY_FIELDS.get(category, [])
    lines = [f"Рубрика: {category}"]
    for field in fields:
        val = values.get(field["key"], "")
        if val:
            lines.append(f"{field['label']}: {val}")
    return "\n".join(lines)

# ============================================================
# SESSION STATE
# ============================================================
defaults = {
    "current_text": "",
    "chat_history": [],
    "raw_source_data": "",
    "versions_history": [],
}
for key, value in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = value

_flags = {
    "show_metrics": False,
    "show_yandex": False,
    "show_news": False,
    "news_for_post": None,
    "disk_content": "",
    "uploader_nonce": 0,
    "posts_buster": 0,
    "news_cache_buster": 0,
    "disk_cache_buster": 0,
    "fact_check_result": "",
    "last_news_count": 0,
    "pending_download": None,
}
for key, value in _flags.items():
    if key not in st.session_state:
        st.session_state[key] = value

# Разовая проверка доступа к Диску, чтобы индикатор в сайдбаре не горел
# зелёным при неверном токене.
if storage is not None and not st.session_state.get("storage_probed"):
    storage.probe()
    st.session_state.storage_probed = True
    st.session_state.storage_writable = storage.writable
    st.session_state.storage_error = storage.last_error
elif storage is not None:
    storage.writable = st.session_state.get("storage_writable")
    storage.last_error = st.session_state.get("storage_error", "")

# ============================================================
# БОКОВАЯ ПАНЕЛЬ
# ============================================================
with st.sidebar:
    render_sidebar_brand()
    render_storage_status()

    if storage is None:
        st.caption("История и примеры сохраняются локально до перезапуска.")
    elif storage.writable is False:
        st.caption(f"Ошибка Яндекс.Диска: {storage.last_error}")

    if not APP_PASSWORD:
        st.warning("Пароль доступа не задан — приложение открыто всем, у кого есть ссылка.")

    st.markdown('<div class="sidebar-section-label">Контент</div>', unsafe_allow_html=True)
    task = st.selectbox("Задача", ["Написать пост для Telegram-канала"])
    post_category = st.selectbox("Рубрика", list(CATEGORY_FIELDS.keys()))

    with st.expander("Параметры публикации", expanded=True):
        text_length = st.select_slider(
            "Объём",
            options=["Короткий (до 1000 зн.)", "Стандартный", "Развернутый (лонгрид)"],
            value="Стандартный",
        )
        tone = st.selectbox(
            "Тональность",
            ["Строгий (сухие факты)", "Стандарт (информационный)", "Живой (для соцсетей)"],
        )
        use_citations = st.checkbox("Использовать цитаты", value=False)

    with st.expander("Примеры стиля", expanded=False):
        st.caption("Добавьте удачные публикации — они будут использоваться как ориентир по подаче.")
        style_files = st.file_uploader(
            "Файлы",
            type=list(STYLE_EXTENSIONS),
            accept_multiple_files=True,
            key="style_files_uploader",
        )
        pasted_examples = st.text_area(
            "Или вставьте текст",
            key="style_paste_area",
            height=80,
        )
        if st.button("Добавить примеры", use_container_width=True):
            new_texts = []
            for f in style_files or []:
                t = extract_text_from_file(f)
                if t.strip():
                    new_texts.append(t)
            if (pasted_examples or "").strip():
                new_texts.extend(
                    [chunk.strip() for chunk in pasted_examples.split("\n---\n") if chunk.strip()]
                )
            if new_texts:
                if add_style_examples(new_texts):
                    st.success(f"Добавлено: {len(new_texts)}")
                    st.rerun()
            else:
                st.warning("Не удалось извлечь ни одного примера.")

        library = [e for e in load_style_library() if e.get("id")]
        if library:
            if "selected_style_ids" not in st.session_state:
                st.session_state.selected_style_ids = [
                    e["id"] for e in library[-MAX_STYLE_EXAMPLES_IN_PROMPT:]
                ]

            st.caption(f"В библиотеке: {len(library)}")
            for e in reversed(library):
                example_id = e["id"]
                preview = (e.get("preview") or e.get("text") or "Без названия")[:50]
                col1, col2 = st.columns([4, 1])
                with col1:
                    checked = st.checkbox(
                        preview,
                        value=example_id in st.session_state.selected_style_ids,
                        key=widget_key("chk", example_id),
                    )
                    if checked and example_id not in st.session_state.selected_style_ids:
                        st.session_state.selected_style_ids.append(example_id)
                    elif not checked and example_id in st.session_state.selected_style_ids:
                        st.session_state.selected_style_ids.remove(example_id)
                with col2:
                    if st.button("×", key=widget_key("del", example_id), use_container_width=True):
                        if delete_style_example(example_id):
                            st.rerun()

    st.markdown('<div class="sidebar-section-label">Навигация</div>', unsafe_allow_html=True)
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        if st.button("↻", help="Новый пост", use_container_width=True):
            reset_input_form()
            st.rerun()
    with col2:
        if st.button("▥", help="Аналитика", use_container_width=True):
            st.session_state.show_metrics = True
            st.session_state.show_yandex = False
            st.session_state.show_news = False
            st.rerun()
    with col3:
        if st.button("◫", help="База знаний", use_container_width=True):
            st.session_state.show_yandex = True
            st.session_state.show_metrics = False
            st.session_state.show_news = False
            st.rerun()
    with col4:
        if st.button("◎", help="Мониторинг", use_container_width=True):
            st.session_state.show_news = True
            st.session_state.show_metrics = False
            st.session_state.show_yandex = False
            st.rerun()

# ============================================================
# ОСНОВНОЙ ЭКРАН
# ============================================================
if not st.session_state.show_metrics and not st.session_state.show_yandex and not st.session_state.show_news:
    render_main_hero()
    left_column, right_column = st.columns([1.06, 0.94], gap="large")

    with left_column:
        render_section_header("01 · INPUT", "Исходные данные", "Фактура, тезисы и материалы для будущей публикации")
        with st.container(border=True):
            st.caption(f"Рубрика: {post_category}")
            field_values = build_dynamic_fields(post_category)

            speaker = st.text_input(
                "👤 Спикер",
                placeholder="Имя, фамилия, должность",
                key="speaker_input",
            )
            materials_text = st.text_area(
                "📝 Материалы / тезисы", height=100, key="materials_input"
            )

            # Новость, выбранная в разделе «Мониторинг» кнопкой 📋
            news_block = ""
            if st.session_state.news_for_post:
                n = st.session_state.news_for_post
                with st.container(border=True):
                    c1, c2 = st.columns([5, 1])
                    with c1:
                        st.caption("📰 Новость как источник")
                        st.markdown(f"**{n.get('title') or 'Без заголовка'}**")
                        st.caption(f"{n.get('source') or '—'} • {format_news_date(n)}")
                    with c2:
                        if st.button("✕", key="drop_news", use_container_width=True):
                            st.session_state.news_for_post = None
                            st.rerun()
                news_block = (
                    f"Новость-источник: {n.get('title') or ''}\n"
                    f"Издание: {n.get('source') or ''}\n"
                    f"Дата: {format_news_date(n)}\n"
                    f"Ссылка: {n.get('url') or ''}\n"
                    f"Краткое содержание: {n.get('description') or ''}"
                )
            
            source_tab1, source_tab2 = st.tabs(["📤 Загрузить файлы", "☁️ Яндекс.Диск"])
            
            uploaded_files = None
            with source_tab1:
                uploaded_files = st.file_uploader(
                    "📁 Файлы (PDF, DOCX, PPTX, XLSX, TXT, HTML)",
                    type=list(SUPPORTED_EXTENSIONS),
                    accept_multiple_files=True,
                    key=f"main_uploader_{st.session_state.uploader_nonce}",
                )
            
            with source_tab2:
                if disk_api:
                    st.info("📂 Файлы из базы знаний ЦСР")
                    base_path = YANDEX_BASE_PATH
                    folders_data = cached_list_disk(base_path, disk_buster())

                    if folders_data:
                        folders = [f for f in folders_data if f.get("type") == "dir"]
                        files = [f for f in folders_data if f.get("type") == "file"]

                        folder_names = ["📁 Все файлы"] + [
                            f"📁 {f.get('name', '')}" for f in folders
                        ]
                        selected_folder_display = st.selectbox(
                            "Выберите тему:",
                            folder_names,
                            key="disk_folder_select"
                        )

                        if selected_folder_display == "📁 Все файлы":
                            files_to_show = files
                        else:
                            folder_name = selected_folder_display.replace("📁 ", "", 1)
                            folder_obj = next(
                                (f for f in folders if f.get("name") == folder_name), None
                            )
                            if folder_obj:
                                files_to_show = [
                                    f for f in cached_list_disk(folder_obj.get("path"), disk_buster())
                                    if f.get("type") == "file"
                                ]
                            else:
                                files_to_show = []

                        readable = [
                            f for f in files_to_show
                            if (f.get("name") or "").lower().endswith(
                                tuple(f".{e}" for e in SUPPORTED_EXTENSIONS)
                            )
                        ]
                        skipped = len(files_to_show) - len(readable)

                        if readable:
                            st.write(f"**Файлы:** {len(readable)}")
                            if skipped:
                                # Раньше .doc и .xls предлагались к выбору,
                                # хотя парсер их не поддерживает, и молча
                                # давали пустой результат.
                                st.caption(
                                    f"Ещё {skipped} файлов пропущено: формат не поддерживается."
                                )
                            selected_disk_files = []

                            for f in readable:
                                name = f.get("name") or ""
                                if st.checkbox(
                                    f"{file_icon(name)} {name}",
                                    key=widget_key("disk", f.get("path")),
                                ):
                                    selected_disk_files.append(f)

                            if selected_disk_files:
                                if st.button("✓ Добавить с диска", use_container_width=True):
                                    disk_parts = []
                                    with st.spinner("Читаю файлы с Диска..."):
                                        for f in selected_disk_files:
                                            content = extract_yandex_file(disk_api, f.get("path"))
                                            if content:
                                                disk_parts.append(
                                                    f"### {f.get('name')}\n{content}"
                                                )
                                    if disk_parts:
                                        st.session_state.disk_content = "\n\n".join(disk_parts)
                                        st.success(f"✅ Загружено файлов: {len(disk_parts)}")
                                    else:
                                        st.warning("Не удалось извлечь текст из выбранных файлов.")
                        else:
                            st.info("📁 В выбранной папке нет поддерживаемых файлов")
                    else:
                        st.warning("❌ База знаний не найдена")
                else:
                    st.warning("⚠️ Яндекс.Диск не подключен")

            if st.button("▶️ Сформировать", type="primary", use_container_width=True):
                has_field_data = any((v or "").strip() for v in field_values.values())
                disk_content = st.session_state.get("disk_content", "")

                if not has_field_data and not (materials_text or "").strip() and not uploaded_files and not disk_content and not news_block:
                    st.warning("⚠️ Добавьте данные")
                else:
                    extracted_file_text = extract_text_from_files(uploaded_files) if uploaded_files else ""
                    fields_text = format_fields_as_text(post_category, field_values)

                    final_materials = "\n".join(filter(None, [
                        fields_text,
                        f"Спикер: {speaker}" if speaker else "",
                        materials_text,
                        news_block,
                        extracted_file_text,
                        disk_content
                    ]))
                    st.session_state.raw_source_data = final_materials

                    style_block = build_style_block(st.session_state.get("selected_style_ids", []))
                    system_rules = get_enhanced_system_rules(post_category, tone, text_length, use_citations)
                    if style_block:
                        system_rules += "\n" + style_block

                    st.session_state.chat_history = [
                        {"role": "system", "content": system_rules},
                        {"role": "user", "content": st.session_state.raw_source_data},
                    ]

                    with st.spinner("✨ Создаю пост..."):
                        result = call_openai(st.session_state.chat_history, TEMP_GENERATE)
                    if result:
                        st.session_state.current_text = result
                        st.session_state.chat_history.append(
                            {"role": "assistant", "content": result}
                        )
                        st.session_state.versions_history = [
                            {"text": result, "label": "v1"}
                        ]
                        st.session_state.fact_check_result = ""
                        save_post(
                            post_category, tone, text_length, result,
                            st.session_state.raw_source_data,
                            1 if use_citations else 0,
                        )
                        st.session_state.disk_content = ""
                        st.rerun()

    with right_column:
        render_section_header("02 · OUTPUT", "Готовый материал", "Редактирование, проверка и экспорт результата")
        with st.container(border=True):
            text_tab, fact_tab, history_tab, export_tab = st.tabs(["📄 Текст", "🔍 Проверка", "📋 История", "📥 Экспорт"])

            with text_tab:
                if st.session_state.current_text:
                    st.code(st.session_state.current_text, language="markdown")

                    st.divider()
                    st.write("🔄 **Правки:**")
                    with st.form("refine_form", clear_on_submit=True):
                        refine_prompt = st.text_input("Что изменить?", placeholder="Пример: сделай короче")
                        submitted = st.form_submit_button("✏️ Уточнить")
                        if submitted and refine_prompt.strip():
                            st.session_state.chat_history.append({"role": "user", "content": refine_prompt})
                            api_messages = trim_history(st.session_state.chat_history)
                            with st.spinner("✏️ Переписываю..."):
                                result = call_openai(api_messages, TEMP_REFINE)
                                if result:
                                    st.session_state.current_text = result
                                    st.session_state.chat_history.append({"role": "assistant", "content": result})
                                    v_num = len(st.session_state.versions_history) + 1
                                    st.session_state.versions_history.append({"text": result, "label": f"v{v_num}"})
                                    st.rerun()
                else:
                    st.info("👈 Заполните данные слева")

            with fact_tab:
                if st.session_state.current_text:
                    if st.button("🔍 Проверить", type="primary", use_container_width=True):
                        fact_prompt = f"""Проверь текст на ошибки:
1. Цифры, даты, фамилии — без искажений
2. Нет дублирующихся блоков
3. Правильность названий

ИСХОДНИКИ: {st.session_state.raw_source_data[:4000]}
ТЕКСТ: {st.session_state.current_text[:4000]}"""
                        with st.spinner("🔍 Проверяю..."):
                            check = call_openai(
                                [{"role": "user", "content": fact_prompt}], TEMP_FACTCHECK
                            )
                        # Раньше результат жил только до следующей перерисовки
                        # и пропадал от любого клика в интерфейсе.
                        st.session_state.fact_check_result = check or ""

                    if st.session_state.fact_check_result:
                        st.success("✅ Готово!")
                        st.write(st.session_state.fact_check_result)
                else:
                    st.info("Сначала создайте текст")

            with history_tab:
                if st.session_state.versions_history:
                    total = len(st.session_state.versions_history)
                    for offset, version in enumerate(reversed(st.session_state.versions_history)):
                        idx = total - offset  # стабильный номер версии
                        col1, col2 = st.columns([4, 1])
                        with col1:
                            st.markdown(f"**{version.get('label', f'v{idx}')}**")
                            st.caption((version.get("text") or "")[:120].replace("\n", " "))
                        with col2:
                            if st.button("✓", key=f"restore_v{idx}", use_container_width=True):
                                st.session_state.current_text = version.get("text", "")
                                st.rerun()
                        st.divider()
                else:
                    st.info("История версий")

            with export_tab:
                if st.session_state.current_text:
                    docx_data = create_docx(st.session_state.current_text)
                    col1, col2, col3 = st.columns(3)
                    stamp = datetime.now().strftime("%Y%m%d_%H%M")
                    with col1:
                        st.download_button(
                            "📄 Word",
                            data=docx_data,
                            file_name=f"post_{stamp}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key="export_docx",
                        )
                    with col2:
                        st.download_button(
                            "📝 Markdown",
                            data=st.session_state.current_text.encode("utf-8"),
                            file_name=f"post_{stamp}.md",
                            mime="text/markdown",
                            key="export_md",
                        )
                    with col3:
                        st.download_button(
                            "📋 TXT",
                            data=st.session_state.current_text.encode("utf-8"),
                            file_name=f"post_{stamp}.txt",
                            mime="text/plain",
                            key="export_txt",
                        )
                else:
                    st.info("Создайте текст")

# ============================================================
# DASHBOARD МЕТРИК
# ============================================================
elif st.session_state.show_metrics:
    render_section_header("DASHBOARD", "Аналитика и метрики", "Динамика подготовки контента и структура публикаций")
    
    if st.button("← Вернуться"):
        st.session_state.show_metrics = False
        st.rerun()
    
    st.divider()
    
    metrics = get_metrics()
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-orb"></div>
            <div class="metric-label">Всего постов</div>
            <div class="metric-value">{metrics['total_posts']}</div>
            <div class="metric-note">Все сохранённые публикации</div>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-orb"></div>
            <div class="metric-label">За неделю</div>
            <div class="metric-value">{metrics['posts_week']}</div>
            <div class="metric-note">Создано за последние 7 дней</div>
        </div>
        """, unsafe_allow_html=True)
    with col3:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-orb"></div>
            <div class="metric-label">За месяц</div>
            <div class="metric-value">{metrics['posts_month']}</div>
            <div class="metric-note">Активность за последние 30 дней</div>
        </div>
        """, unsafe_allow_html=True)
    with col4:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-orb"></div>
            <div class="metric-label">С цитатами</div>
            <div class="metric-value">{metrics['with_citations']}</div>
            <div class="metric-note">Материалы с прямой речью</div>
        </div>
        """, unsafe_allow_html=True)
    
    st.divider()
    
    col1, col2 = st.columns(2)
    
    def chart_from_pairs(pairs, label, strip_numbers=False):
        """Строит график с ЧИТАЕМЫМИ подписями по оси X."""
        if not pairs:
            st.info("Нет данных")
            return
        names = []
        for name, _ in pairs:
            n = name or "—"
            if strip_numbers:
                n = re.sub(r"^\d+\.\s*", "", n)
            names.append(n[:28])
        counts = [c for _, c in pairs]
        df = pd.DataFrame({label: counts}, index=names)
        st.bar_chart(df)

    with col1:
        st.subheader("📂 По рубрикам")
        chart_from_pairs(metrics['by_category'], "Постов", strip_numbers=True)

    with col2:
        st.subheader("🎙️ По тону")
        chart_from_pairs(metrics['by_tone'], "Постов")

    st.divider()

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("📏 По объему")
        chart_from_pairs(metrics['by_length'], "Постов")

    with col2:
        st.subheader("💬 Использование цитат")
        # Раньше здесь передавался словарь из чисел — Streamlit падал с ошибкой
        chart_from_pairs(
            [("С цитатами", metrics['with_citations']),
             ("Без цитат", metrics['without_citations'])],
            "Постов"
        )
    
    st.divider()
    
    st.subheader("📋 История последних постов")
    posts = get_posts_history(10)
    if posts:
        for post_id, category, tone, text, created_at, has_citations in posts:
            text = text or ""
            label = f"{category or '—'} • {(created_at or '')[:10] or 'без даты'}"
            if has_citations:
                label += " • 💬"
            with st.expander(label):
                st.write(text[:200] + ("..." if len(text) > 200 else ""))
    else:
        st.info("Нет постов")

# ============================================================
# ЯНДЕКС.ДИСК БРАУЗЕР
# ============================================================
elif st.session_state.show_yandex:
    render_section_header("KNOWLEDGE BASE", "База знаний ЦСР", "Файлы и тематические материалы на Яндекс.Диске")
    
    if st.button("← Вернуться"):
        st.session_state.show_yandex = False
        st.rerun()
    
    st.divider()
    
    if not disk_api:
        st.error("⚠️ Яндекс.Диск не подключен")
        st.info("Добавьте YANDEX_DISK_TOKEN в Streamlit Cloud secrets")
    else:
        base_path = YANDEX_BASE_PATH
        
        col1, col2 = st.columns([3, 1])
        with col1:
            st.write(f"**📂 {base_path}**")
        with col2:
            # Раньше кнопка просто перезапускала скрипт, а листинг ещё
            # 5 минут отдавался из кэша.
            if st.button("🔄 Обновить", use_container_width=True):
                bump_disk_cache()
                st.rerun()

        st.divider()

        items = cached_list_disk(base_path, disk_buster())

        if not items:
            st.error(f"Папка «{base_path}» не найдена или пуста.")
            st.caption(
                "Имя папки на Диске должно совпадать с константой YANDEX_BASE_PATH "
                "в начале файла. Ниже — что реально лежит в корне вашего Диска:"
            )
            root = cached_list_disk("/", disk_buster())
            if root:
                for r in root:
                    kind = "📁" if r.get("type") == "dir" else "📄"
                    st.code(f'{kind}  {r.get("path")}')
                st.caption(
                    "Скопируйте нужный путь и подставьте его в YANDEX_BASE_PATH."
                )
            else:
                st.warning(
                    "Корень Диска тоже не читается — скорее всего, неверный "
                    "YANDEX_DISK_TOKEN или у токена нет прав на чтение Диска."
                )

        if items:
            folders = [f for f in items if f.get("type") == "dir"]
            files = [f for f in items if f.get("type") == "file"]
            
            st.write(f"**📊 Статистика:** {len(folders)} папок, {len(files)} файлов")
            st.divider()
            
            if folders:
                st.subheader("📁 Тематические папки")
                for folder in sorted(folders, key=lambda x: x.get("name") or ""):
                    folder_name = folder.get("name") or ""
                    with st.expander(f"📁 {folder_name}"):
                        sub_items = cached_list_disk(folder.get("path"), disk_buster())
                        sub_files = [f for f in sub_items if f.get("type") == "file"]

                        if sub_files:
                            st.write(f"**{len(sub_files)} файлов:**")
                            for f in sorted(sub_files, key=lambda x: x.get("name") or ""):
                                render_disk_file_row(f, "dl")
                        else:
                            st.info("Папка пуста")

            if files:
                st.divider()
                st.subheader("📄 Файлы в корне")
                for f in sorted(files, key=lambda x: x.get("name") or ""):
                    render_disk_file_row(f, "dl_root")

# ============================================================
# МОНИТОРИНГ НОВОСТЕЙ
# ============================================================
elif st.session_state.show_news:
    render_section_header("MEDIA RADAR", "Мониторинг новостей ЦСР", "Поиск инфоповодов и сохранение источников для публикаций")
    
    if st.button("← Вернуться"):
        st.session_state.show_news = False
        st.rerun()
    
    st.divider()
    
    col1, col2, col3 = st.columns([2, 1, 1])
    
    with col1:
        query = st.text_input(
            "🔍 Поиск:",
            value='ЦСР OR "Центр стратегических разработок"',
            placeholder="Поисковый запрос",
            key="news_query",
        )

    with col2:
        days = st.selectbox("📅 Период:", [7, 14, 30], index=0, key="news_days")

    with col3:
        st.write("")
        if st.button("🔄 Обновить", use_container_width=True):
            st.session_state.news_cache_buster += 1
            st.rerun()

    st.divider()

    tab1, tab2, tab3 = st.tabs(["🔍 Новости", "⭐ Избранное", "📊 Статистика"])

    with tab1:
        st.write("**Последние новости о ЦСР**")

        col_google, col_api = st.columns(2)
        with col_google:
            use_google = st.checkbox("📰 Google News (бесплатно)", value=True, key="use_google")
        with col_api:
            use_api = st.checkbox(
                "🔐 News API",
                value=NEWS_API_KEY is not None,
                disabled=NEWS_API_KEY is None,
                key="use_news_api",
            )

        if not use_google and not use_api:
            st.warning("Выберите хотя бы один источник новостей.")
            news_list = []
        else:
            with st.spinner("📰 Загружаю новости..."):
                news_list = cached_fetch_news(
                    query, use_api, use_google, days,
                    st.session_state.news_cache_buster,
                )

        # Статистика на третьей вкладке раньше опиралась на 'news_list' in
        # locals() — теперь значение живёт в session_state явно.
        st.session_state.last_news_count = len(news_list)

        if news_list:
            st.write(f"**Найдено: {len(news_list)} новостей**")
            st.divider()

            for i, news in enumerate(news_list):
                title = news.get("title") or "Без заголовка"
                source = news.get("source") or "Источник неизвестен"
                published = news.get("published") or ""
                description = news.get("description") or ""
                url = news.get("url") or ""

                with st.container(border=True):
                    col_body, col_actions = st.columns([5, 1])

                    with col_body:
                        if url:
                            st.markdown(f"**[{title}]({url})**")
                        else:
                            st.markdown(f"**{title}**")
                        st.caption(f"📰 {source} • {format_news_date(news)}")

                        if description:
                            st.write(
                                description[:300] + "..."
                                if len(description) > 300 else description
                            )

                    with col_actions:
                        col_a, col_b = st.columns(2)
                        # Ключи по url, а не по индексу: список между
                        # перерисовками может сдвинуться.
                        with col_a:
                            if st.button(
                                "⭐",
                                key=widget_key("fav", url or title, i),
                                help="В избранное",
                                use_container_width=True,
                            ):
                                if news_monitor.save_favorite(
                                    title, source, url, published, description
                                ):
                                    st.toast("Добавлено в избранное")
                                    st.rerun()

                        with col_b:
                            if st.button(
                                "📋",
                                key=widget_key("use", url or title, i),
                                help="В пост",
                                use_container_width=True,
                            ):
                                st.session_state.news_for_post = {
                                    "title": title,
                                    "source": source,
                                    "url": url,
                                    "description": description,
                                    "published": published,
                                    "published_dt": news.get("published_dt"),
                                }
                                st.session_state.show_news = False
                                st.toast("Новость добавлена как источник")
                                # Без rerun пользователь оставался на экране
                                # мониторинга до следующего клика.
                                st.rerun()

                if i < len(news_list) - 1:
                    st.divider()
        else:
            st.info("📭 Новостей не найдено. Попробуйте другой запрос.")

    with tab2:
        st.write("**Избранные новости**")

        favorites = news_monitor.get_favorites()

        if favorites:
            st.write(f"**Всего: {len(favorites)} новостей**")
            st.divider()

            for idx, fav in enumerate(favorites):
                title = fav["title"] or "Без заголовка"
                description = fav["description"]
                with st.container(border=True):
                    col_body, col_actions = st.columns([5, 1])

                    with col_body:
                        if fav["url"]:
                            st.markdown(f"**[{title}]({fav['url']})**")
                        else:
                            st.markdown(f"**{title}**")
                        st.caption(
                            f"📰 {fav['source'] or '—'} • {format_news_date(fav['published'])}"
                        )

                        if description:
                            st.write(
                                description[:200] + "..."
                                if len(description) > 200 else description
                            )

                    with col_actions:
                        # Прежний ключ обрезал заголовок до 20 символов —
                        # однотипные заголовки давали DuplicateElementKey.
                        if st.button(
                            "🗑️",
                            key=widget_key("del_fav", title, idx),
                            use_container_width=True,
                        ):
                            if news_monitor.delete_favorite(title):
                                st.rerun()

                st.divider()
        else:
            st.info("⭐ Избранные новости пока пусты")
    
    with tab3:
        st.write("**Статистика мониторинга**")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-orb"></div>
                <div class="metric-label">Найдено новостей</div>
                <div class="metric-value">{st.session_state.last_news_count}</div>
                <div class="metric-note">Текущий поисковый срез</div>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            favorites = news_monitor.get_favorites()
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-orb"></div>
                <div class="metric-label">В избранном</div>
                <div class="metric-value">{len(favorites)}</div>
                <div class="metric-note">Сохранённые инфоповоды</div>
            </div>
            """, unsafe_allow_html=True)
        
        with col3:
            status = "🟢" if NEWS_API_KEY else "🔴"
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-orb"></div>
                <div class="metric-label">News API</div>
                <div class="metric-value">{status}</div>
                <div class="metric-note">Статус внешнего источника</div>
            </div>
            """, unsafe_allow_html=True)
        
        st.divider()
        st.info("""
        📰 **Источники новостей:**
        - **News API** - полнотекстовый поиск (если подключен)
        - **Google News** - RSS парсер (всегда доступен)
        
        💡 **Как использовать:**
        1. Введите поисковый запрос
        2. Выберите источники
        3. Нажмите "Обновить" для поиска
        4. ⭐ Сохраняйте интересные новости
        5. 📋 Добавляйте источники в пост
        """)